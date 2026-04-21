#!/usr/bin/env python3
"""
doc-to-md.py — Convert binary/web documents to Markdown for AI SDLC ingestion.

Supported: .docx  .xlsx  .pptx  .html / .htm  .pdf
Output: <basename>.extracted.md  (written to .sdlc/import/ or alongside file)

Usage:
    python3 scripts/doc-to-md.py <file_or_dir> [--output-dir <dir>]

Requires (install what you need):
    pip install pypdf pdfplumber python-docx openpyxl python-pptx \
                beautifulsoup4 html2text trafilatura

All dependencies are optional per format — the script degrades gracefully if
one library is missing and tries fallback strategies before giving up.
"""

import argparse
import os
import sys
import re
from pathlib import Path

# ── Output dir default ─────────────────────────────────────────────────────────
DEFAULT_OUTPUT_DIR = ".sdlc/import"
MAX_BYTES = 10 * 1024 * 1024  # 10 MB safety cap per file

# ── Helpers ────────────────────────────────────────────────────────────────────

def _try_import(module: str):
    try:
        import importlib
        return importlib.import_module(module)
    except ImportError:
        return None


def sanitize(text: str) -> str:
    """Collapse excessive blank lines."""
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


# ── Format extractors ──────────────────────────────────────────────────────────

def extract_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")

    # Try trafilatura first (best main-content extraction)
    trf = _try_import("trafilatura")
    if trf:
        result = trf.extract(raw, include_tables=True, include_formatting=True)
        if result:
            return sanitize(result)

    # Fallback: html2text
    h2t = _try_import("html2text")
    if h2t:
        converter = h2t.HTML2Text()
        converter.ignore_links = False
        converter.body_width = 0
        return sanitize(converter.handle(raw))

    # Last resort: strip tags with BeautifulSoup
    bs4 = _try_import("bs4")
    if bs4:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return sanitize(soup.get_text(separator="\n"))

    raise RuntimeError("No HTML parser available. Install beautifulsoup4 or html2text.")


def extract_pdf(path: Path) -> str:
    # Try pdfplumber (better table support)
    plumber = _try_import("pdfplumber")
    if plumber:
        lines = []
        with plumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    lines.append(f"<!-- page {i} -->\n{text}")
        if lines:
            return sanitize("\n\n".join(lines))

    # Fallback: pypdf
    pypdf = _try_import("pypdf")
    if pypdf:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"<!-- page {i} -->\n{text}")
        return sanitize("\n\n".join(parts))

    raise RuntimeError("No PDF parser available. Install pypdf or pdfplumber.")


def extract_docx(path: Path) -> str:
    # Try mammoth (best Markdown/HTML output)
    mammoth = _try_import("mammoth")
    if mammoth:
        with open(str(path), "rb") as f:
            result = mammoth.convert_to_markdown(f)
        return sanitize(result.value)

    # Fallback: python-docx paragraph text
    docx = _try_import("docx")
    if docx:
        from docx import Document
        doc = Document(str(path))
        lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else ""
                if style.startswith("Heading"):
                    level = re.sub(r"\D", "", style) or "1"
                    lines.append(f"{'#' * int(level)} {para.text}")
                else:
                    lines.append(para.text)
        return sanitize("\n\n".join(lines))

    raise RuntimeError("No DOCX parser available. Install mammoth or python-docx.")


def _col_letter(n: int) -> str:
    """Convert 0-based column index to A, B, …, AA, AB, …"""
    result = ""
    n += 1
    while n:
        result = chr(65 + (n - 1) % 26) + result
        n = (n - 1) // 26
    return result


def extract_xlsx(path: Path) -> str:
    openpyxl = _try_import("openpyxl")
    if not openpyxl:
        pandas = _try_import("pandas")
        if pandas:
            import pandas as pd
            sheets = pd.read_excel(str(path), sheet_name=None)
            parts = []
            for name, df in sheets.items():
                parts.append(f"## Sheet: {name}\n\n{df.to_markdown(index=False)}")
            return sanitize("\n\n".join(parts))
        raise RuntimeError("No Excel parser available. Install openpyxl or pandas.")

    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # Build Markdown table
        headers = [str(c) if c is not None else "" for c in rows[0]]
        md_rows = ["| " + " | ".join(headers) + " |",
                   "| " + " | ".join(["---"] * len(headers)) + " |"]
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            # Pad if shorter
            while len(cells) < len(headers):
                cells.append("")
            md_rows.append("| " + " | ".join(cells) + " |")
        parts.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(md_rows))
    wb.close()
    return sanitize("\n\n".join(parts))


def extract_pptx(path: Path) -> str:
    pptx = _try_import("pptx")
    if not pptx:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for j, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                if j == 0 and not title:
                    title = text
                else:
                    level = para.level if para.level else 0
                    prefix = "  " * level + "-"
                    bullets.append(f"{prefix} {text}")
        heading = f"## Slide {i}" + (f": {title}" if title else "")
        body = "\n".join(bullets)
        slides.append(f"{heading}\n\n{body}" if body else heading)
    return sanitize("\n\n".join(slides))


# ── Dispatch ───────────────────────────────────────────────────────────────────

EXTRACTORS = {
    ".html": extract_html,
    ".htm":  extract_html,
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    ".pptx": extract_pptx,
}


def convert_file(src: Path, out_dir: Path) -> Path:
    suffix = src.suffix.lower()
    if suffix not in EXTRACTORS:
        raise ValueError(f"Unsupported format: {suffix}  (supported: {', '.join(EXTRACTORS)})")

    if src.stat().st_size > MAX_BYTES:
        raise ValueError(f"File too large (>{MAX_BYTES // 1_048_576} MB): {src}")

    print(f"  extracting {src.name} ...", end=" ", flush=True)
    md = EXTRACTORS[suffix](src)
    print(f"{len(md):,} chars")

    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / (src.stem + ".extracted.md")

    # Prepend provenance header
    header = (
        f"<!-- source: {src.name}  converted-by: doc-to-md.py -->\n"
        f"<!-- original-path: {src.resolve()} -->\n\n"
    )
    out_path.write_text(header + md, encoding="utf-8")
    return out_path


# ── CLI ────────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Convert Office / PDF / HTML files to Markdown for AI SDLC ingestion."
    )
    parser.add_argument("inputs", nargs="+",
                        help="File(s) or directory to convert")
    parser.add_argument("--output-dir", "-o", default=None,
                        help=f"Output directory (default: {DEFAULT_OUTPUT_DIR} relative to cwd, "
                             "or alongside each source file)")
    args = parser.parse_args()

    errors = []
    converted = []

    for raw_path in args.inputs:
        src = Path(raw_path)
        if src.is_dir():
            files = [f for f in src.rglob("*") if f.suffix.lower() in EXTRACTORS]
        elif src.is_file():
            files = [src]
        else:
            print(f"  SKIP {raw_path} — not found", file=sys.stderr)
            continue

        for f in files:
            out_dir = Path(args.output_dir) if args.output_dir else Path(DEFAULT_OUTPUT_DIR)
            try:
                out = convert_file(f, out_dir)
                converted.append(out)
                print(f"  -> {out}")
            except Exception as exc:
                errors.append((f, str(exc)))
                print(f"  ERROR {f}: {exc}", file=sys.stderr)

    print(f"\nDone. {len(converted)} converted, {len(errors)} errors.")
    if errors:
        sys.exit(1)


if __name__ == "__main__":
    main()
